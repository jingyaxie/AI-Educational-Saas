from django.shortcuts import render
from rest_framework.views import APIView
from rest_framework.response import Response
from rest_framework import status, permissions
from django.contrib.auth import authenticate, login
from .models import User, UserGroup, Agent, ModelApi, TokenUsage, KnowledgeBase, Space, SpaceMember, SpaceDocument, KnowledgeFile
from .serializers import UserSerializer, UserGroupSerializer, AgentSerializer, ModelApiSerializer, TokenUsageSerializer, KnowledgeBaseSerializer, SpaceSerializer, SpaceMemberSerializer, SpaceDocumentSerializer, KnowledgeFileSerializer
from django.contrib.auth.decorators import login_required
from rest_framework.permissions import IsAuthenticated
from rest_framework_simplejwt.tokens import RefreshToken
from django.views.decorators.csrf import csrf_exempt
from django.utils.decorators import method_decorator
from captcha.models import CaptchaStore
from captcha.helpers import captcha_image_url
from rest_framework.permissions import AllowAny
from django.urls import reverse
from rest_framework import generics, filters
from django_filters.rest_framework import DjangoFilterBackend
import logging
import requests
from django.conf import settings
from rest_framework.decorators import api_view, permission_classes
from rest_framework.parsers import MultiPartParser, FormParser
from rest_framework import serializers
import traceback
import os
import pandas as pd
from docx import Document
import pdfplumber
from pptx import Presentation
from ebooklib import epub
from bs4 import BeautifulSoup
import markdown

from langchain.text_splitter import RecursiveCharacterTextSplitter, CharacterTextSplitter, TokenTextSplitter
from langchain_community.embeddings import OpenAIEmbeddings, HuggingFaceEmbeddings
from langchain_community.vectorstores import Chroma
import openai
from sentence_transformers import SentenceTransformer
openai.api_key = getattr(settings, "OPENAI_API_KEY", os.environ.get("OPENAI_API_KEY"))

logger = logging.getLogger(__name__)

# Create your views here.

class CaptchaView(APIView):
    permission_classes = [AllowAny]
    def get(self, request):
        try:
            print('开始生成验证码')
            new_captcha = CaptchaStore.generate_key()
            print(f'生成验证码 key: {new_captcha}')
            image_url = request.build_absolute_uri(reverse('captcha-image', kwargs={'key': new_captcha}))
            print(f'生成验证码图片 URL: {image_url}')
            return Response({
                'captcha_key': new_captcha,
                'image_url': image_url
            })
        except Exception as e:
            print(f'生成验证码时发生错误: {str(e)}')
            return Response({
                'error': '生成验证码失败',
                'detail': str(e)
            }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

class LoginView(APIView):
    permission_classes = [AllowAny]
    def post(self, request):
        print('收到登录请求')
        print('请求方法:', request.method)
        print('请求头:', request.headers)
        print('请求数据:', request.data)
        username = request.data.get('username')
        password = request.data.get('password')
        captcha_key = request.data.get('captcha_key')
        captcha_value = request.data.get('captcha_value')
        print(f'登录请求：username={username}, captcha_key={captcha_key}, captcha_value={captcha_value}')
        # 校验验证码
        print('开始验证码校验')
        if not CaptchaStore.objects.filter(hashkey=captcha_key, response=captcha_value).exists():
            print('验证码错误')
            return Response({
                'error': '验证码错误',
                'detail': '请输入正确的验证码'
            }, status=status.HTTP_400_BAD_REQUEST)
        print('验证码正确，开始验证用户名密码')
        user = authenticate(request, username=username, password=password)
        if user is not None:
            print(f'用户 {username} 登录成功')
            refresh = RefreshToken.for_user(user)
            print('生成 token 成功')
            return Response({
                'refresh': str(refresh),
                'access': str(refresh.access_token),
                'user': UserSerializer(user).data
            })
        else:
            print(f'用户 {username} 登录失败：用户名或密码错误')
            return Response({'error': '用户名或密码错误'}, status=status.HTTP_401_UNAUTHORIZED)

class UserInfoView(APIView):
    permission_classes = [IsAuthenticated]
    def get(self, request):
        user = request.user
        return Response(UserSerializer(user).data)

class SetRoleView(APIView):
    permission_classes = [IsAuthenticated]
    def post(self, request):
        if request.user.role != 'admin':
            return Response({'error': '无权限'}, status=status.HTTP_403_FORBIDDEN)
        user_id = request.data.get('user_id')
        role = request.data.get('role')
        try:
            user = User.objects.get(id=user_id)
            user.role = role
            user.save()
            return Response({'msg': '角色更新成功'})
        except User.DoesNotExist:
            return Response({'error': '用户不存在'}, status=status.HTTP_404_NOT_FOUND)

@login_required
def dashboard(request):
    context = {
        'username': request.user.username,
        'student_count': 78,
        'graduate_count': 32,
        'token_usage': 109887,
        'knowledge_count': 998,
        'cloud_space_count': 876987,
    }
    return render(request, 'dashboard.html', context)

class MemberListView(generics.ListCreateAPIView):
    queryset = User.objects.all().order_by('-date_joined')
    serializer_class = UserSerializer
    permission_classes = [IsAuthenticated]
    filter_backends = [DjangoFilterBackend, filters.SearchFilter, filters.OrderingFilter]
    filterset_fields = ['group', 'role']
    search_fields = ['username']
    ordering_fields = ['date_joined', 'valid_until']

    def get_queryset(self):
        queryset = super().get_queryset()
        reg_start = self.request.query_params.get('reg_start')
        reg_end = self.request.query_params.get('reg_end')
        if reg_start:
            queryset = queryset.filter(date_joined__gte=reg_start)
        if reg_end:
            queryset = queryset.filter(date_joined__lte=reg_end)
        return queryset

class MemberDetailView(generics.RetrieveUpdateDestroyAPIView):
    queryset = User.objects.all()
    serializer_class = UserSerializer
    permission_classes = [IsAuthenticated]
    lookup_field = 'id'

class UserGroupListView(generics.ListCreateAPIView):
    queryset = UserGroup.objects.all().order_by('-created_at')
    serializer_class = UserGroupSerializer
    permission_classes = [IsAuthenticated]

class UserGroupDetailView(generics.RetrieveUpdateDestroyAPIView):
    queryset = UserGroup.objects.all()
    serializer_class = UserGroupSerializer
    permission_classes = [IsAuthenticated]
    lookup_field = 'id'

class AgentListCreateView(generics.ListCreateAPIView):
    """智能体列表与创建接口"""
    queryset = Agent.objects.all().order_by('-created')
    serializer_class = AgentSerializer
    permission_classes = [IsAuthenticated]

    def get(self, request, *args, **kwargs):
        logger.info('获取智能体列表')
        return super().get(request, *args, **kwargs)

    def post(self, request, *args, **kwargs):
        logger.info(f'创建智能体，请求数据: {request.data}')
        return super().post(request, *args, **kwargs)

class AgentRetrieveUpdateDestroyView(generics.RetrieveUpdateDestroyAPIView):
    """智能体详情、更新、删除接口"""
    queryset = Agent.objects.all()
    serializer_class = AgentSerializer
    permission_classes = [IsAuthenticated]
    lookup_field = 'id'

    def get(self, request, *args, **kwargs):
        logger.info(f'获取智能体详情: {kwargs.get("id")}')
        return super().get(request, *args, **kwargs)

    def put(self, request, *args, **kwargs):
        logger.info(f'更新智能体: {kwargs.get("id")}, 数据: {request.data}')
        return super().put(request, *args, **kwargs)

    def delete(self, request, *args, **kwargs):
        logger.info(f'删除智能体: {kwargs.get("id")}')
        return super().delete(request, *args, **kwargs)

class ModelApiListCreateView(generics.ListCreateAPIView):
    """大模型API列表与创建接口"""
    queryset = ModelApi.objects.all().order_by('-time')
    serializer_class = ModelApiSerializer
    permission_classes = [IsAuthenticated]

    def get(self, request, *args, **kwargs):
        logger.info('获取大模型API列表')
        return super().get(request, *args, **kwargs)

    def post(self, request, *args, **kwargs):
        logger.info(f'创建大模型API，请求数据: {request.data}')
        serializer = self.get_serializer(data=request.data)
        if serializer.is_valid():
            self.perform_create(serializer)
            logger.info('创建大模型API成功')
            return Response(serializer.data, status=status.HTTP_201_CREATED)
        else:
            logger.error(f'创建大模型API失败，校验错误: {serializer.errors}')
            return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

class ModelApiRetrieveUpdateDestroyView(generics.RetrieveUpdateDestroyAPIView):
    """大模型API详情、更新、删除接口"""
    queryset = ModelApi.objects.all()
    serializer_class = ModelApiSerializer
    permission_classes = [IsAuthenticated]
    lookup_field = 'id'

    def get(self, request, *args, **kwargs):
        logger.info(f'获取大模型API详情: {kwargs.get("id")}')
        return super().get(request, *args, **kwargs)

    def put(self, request, *args, **kwargs):
        logger.info(f'更新大模型API: {kwargs.get("id")}, 数据: {request.data}')
        return super().put(request, *args, **kwargs)

    def delete(self, request, *args, **kwargs):
        logger.info(f'删除大模型API: {kwargs.get("id")}')
        return super().delete(request, *args, **kwargs)

@api_view(['POST'])
@permission_classes([IsAuthenticated])
def refresh_usage(request):
    """
    批量刷新所有大模型API的token用量（以OpenAI为例）。
    需在settings中配置 OPENAI_ORG_ID 和 OPENAI_MANAGE_KEY（有用量查询权限的key）。
    """
    from .models import ModelApi
    import datetime
    logger.info('开始批量刷新大模型API用量...')
    results = []
    for api in ModelApi.objects.all():
        usage = None
        if api.model == 'openai':
            try:
                # 查询本月用量（单位：美元，token数需换算）
                org_id = getattr(settings, 'OPENAI_ORG_ID', None)
                manage_key = getattr(settings, 'OPENAI_MANAGE_KEY', None)
                if not org_id or not manage_key:
                    logger.warning('未配置OPENAI_ORG_ID或OPENAI_MANAGE_KEY，无法查询OpenAI用量')
                    continue
                now = datetime.datetime.utcnow()
                start_date = now.replace(day=1).strftime('%Y-%m-%d')
                end_date = now.strftime('%Y-%m-%d')
                url = f'https://api.openai.com/v1/dashboard/billing/usage?start_date={start_date}&end_date={end_date}'
                headers = {
                    'Authorization': f'Bearer {api.apikey}',
                    'OpenAI-Organization': org_id
                }
                resp = requests.get(url, headers=headers, timeout=10)
                logger.info(f'OpenAI用量查询响应: {resp.status_code} {resp.text}')
                if resp.status_code == 200:
                    data = resp.json()
                    usage = data.get('total_usage', 0)
                    # OpenAI返回单位为美分，token数需结合价格换算，这里直接显示美元
                    api.usage = f"${usage/100:.2f}"
                    api.save()
                else:
                    logger.warning(f'OpenAI用量查询失败: {resp.text}')
            except Exception as e:
                logger.error(f'查询OpenAI用量异常: {e}')
        # 其它大模型可在此扩展
        results.append({
            'id': api.id,
            'model': api.model,
            'apikey': api.apikey[:8] + '...' if api.apikey else '',
            'usage': api.usage
        })
    logger.info('批量刷新大模型API用量完成')
    return Response({'results': results}, status=status.HTTP_200_OK)

class TokenUsageListCreateView(generics.ListCreateAPIView):
    queryset = TokenUsage.objects.all().order_by('-created')
    serializer_class = TokenUsageSerializer
    permission_classes = [IsAuthenticated]
    filter_backends = [DjangoFilterBackend, filters.SearchFilter, filters.OrderingFilter]
    filterset_fields = ['user', 'agent', 'apikey']
    search_fields = ['prompt', 'user__username']
    ordering_fields = ['created', 'tokens']

    def get_queryset(self):
        logger.info('[TokenUsage] ====== 开始处理Token消耗明细请求 ======')
        logger.info(f'[TokenUsage] 请求方法: {self.request.method}')
        logger.info(f'[TokenUsage] 请求路径: {self.request.path}')
        logger.info(f'[TokenUsage] 请求参数: {self.request.query_params}')
        logger.info(f'[TokenUsage] 认证信息: {self.request.auth}')
        logger.info(f'[TokenUsage] 当前用户: {self.request.user}')
        
        queryset = super().get_queryset()
        start_date = self.request.query_params.get('start_date')
        end_date = self.request.query_params.get('end_date')
        agent = self.request.query_params.get('agent')
        apikey = self.request.query_params.get('apikey')
        search = self.request.query_params.get('search')
        username = self.request.query_params.get('username')
        user = self.request.query_params.get('user')
        
        logger.info(f'[TokenUsage] 时间范围: {start_date} 至 {end_date}')
        logger.info(f'[TokenUsage] 智能体过滤: {agent}')
        logger.info(f'[TokenUsage] API-key过滤: {apikey}')
        logger.info(f'[TokenUsage] 提示词搜索: {search}')
        logger.info(f'[TokenUsage] 用户名搜索: {username}')
        logger.info(f'[TokenUsage] 用户过滤: {user}')
        logger.info(f'[TokenUsage] 初始查询SQL: {str(queryset.query)}')
        logger.info(f'[TokenUsage] 初始记录数: {queryset.count()}')
        
        # 添加用户过滤
        if user:
            queryset = queryset.filter(user_id=user)
        else:
            queryset = queryset.filter(user=self.request.user)
        logger.info(f'[TokenUsage] 用户过滤后SQL: {str(queryset.query)}')
        logger.info(f'[TokenUsage] 用户过滤后记录数: {queryset.count()}')
        
        # 添加智能体过滤
        if agent:
            queryset = queryset.filter(agent_id=agent)
            logger.info(f'[TokenUsage] 智能体过滤后SQL: {str(queryset.query)}')
            logger.info(f'[TokenUsage] 智能体过滤后记录数: {queryset.count()}')
        
        # 添加API-key过滤
        if apikey:
            queryset = queryset.filter(apikey=apikey)
            logger.info(f'[TokenUsage] API-key过滤后SQL: {str(queryset.query)}')
            logger.info(f'[TokenUsage] API-key过滤后记录数: {queryset.count()}')
        
        # 添加提示词搜索
        if search:
            queryset = queryset.filter(prompt__icontains=search)
            logger.info(f'[TokenUsage] 提示词搜索后SQL: {str(queryset.query)}')
            logger.info(f'[TokenUsage] 提示词搜索后记录数: {queryset.count()}')
        
        # 添加用户名搜索
        if username:
            queryset = queryset.filter(user__username__icontains=username)
            logger.info(f'[TokenUsage] 用户名搜索后SQL: {str(queryset.query)}')
            logger.info(f'[TokenUsage] 用户名搜索后记录数: {queryset.count()}')
        
        # 添加时间范围过滤
        if start_date:
            queryset = queryset.filter(created__gte=start_date)
            logger.info(f'[TokenUsage] 开始日期过滤后SQL: {str(queryset.query)}')
            logger.info(f'[TokenUsage] 开始日期过滤后记录数: {queryset.count()}')
            
        if end_date:
            queryset = queryset.filter(created__lte=end_date)
            logger.info(f'[TokenUsage] 结束日期过滤后SQL: {str(queryset.query)}')
            logger.info(f'[TokenUsage] 结束日期过滤后记录数: {queryset.count()}')
        
        # 记录最终结果
        final_count = queryset.count()
        logger.info(f'[TokenUsage] 最终SQL: {str(queryset.query)}')
        logger.info(f'[TokenUsage] 最终记录数: {final_count}')
        if final_count > 0:
            logger.info(f'[TokenUsage] 返回数据示例: {[str(obj) for obj in queryset[:5]]}')
        else:
            logger.info('[TokenUsage] 没有找到符合条件的记录')
            
        logger.info('[TokenUsage] ====== Token消耗明细请求处理完成 ======')
        return queryset

    def get_serializer_context(self):
        context = super().get_serializer_context()
        context['request'] = self.request
        return context

    def perform_create(self, serializer):
        try:
            logger.info('[TokenUsage] ====== 开始创建Token使用记录 ======')
            logger.info(f'[TokenUsage] 请求数据: {self.request.data}')
            logger.info(f'[TokenUsage] 当前用户: {self.request.user}')
            
            instance = serializer.save()
            
            logger.info(f'[TokenUsage] Token使用记录创建成功: {instance}')
            logger.info('[TokenUsage] ====== Token使用记录创建完成 ======')
        except Exception as e:
            logger.error(f'[TokenUsage] 创建token使用记录失败: {str(e)}')
            logger.error(f'[TokenUsage] 错误详情: {e.__class__.__name__}')
            logger.error(f'[TokenUsage] 错误堆栈: {traceback.format_exc()}')
            raise

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def token_usage_stats(request):
    """获取token使用统计信息"""
    from django.db.models import Sum
    from django.utils import timezone
    import datetime
    
    # 获取时间范围
    today = timezone.now().date()
    yesterday = today - datetime.timedelta(days=1)
    week_start = today - datetime.timedelta(days=today.weekday())
    month_start = today.replace(day=1)
    
    # 计算各时间段的使用量
    total = TokenUsage.objects.aggregate(total=Sum('tokens'))['total'] or 0
    today_usage = TokenUsage.objects.filter(created__date=today).aggregate(total=Sum('tokens'))['total'] or 0
    yesterday_usage = TokenUsage.objects.filter(created__date=yesterday).aggregate(total=Sum('tokens'))['total'] or 0
    week_usage = TokenUsage.objects.filter(created__date__gte=week_start).aggregate(total=Sum('tokens'))['total'] or 0
    month_usage = TokenUsage.objects.filter(created__date__gte=month_start).aggregate(total=Sum('tokens'))['total'] or 0
    
    return Response({
        'total': total,
        'today': today_usage,
        'yesterday': yesterday_usage,
        'week': week_usage,
        'month': month_usage
    })

class KnowledgeBaseListCreateView(generics.ListCreateAPIView):
    queryset = KnowledgeBase.objects.all().order_by('-created')
    serializer_class = KnowledgeBaseSerializer
    permission_classes = [IsAuthenticated]
    filter_backends = [DjangoFilterBackend, filters.SearchFilter, filters.OrderingFilter]
    search_fields = ['name']
    ordering_fields = ['created']

    def create(self, request, *args, **kwargs):
        serializer = self.get_serializer(data=request.data)
        if not serializer.is_valid():
            logger.error(f'KnowledgeBase create 400: {serializer.errors}')
            return Response(serializer.errors, status=400)
        self.perform_create(serializer)
        headers = self.get_success_headers(serializer.data)
        return Response(serializer.data, status=201, headers=headers)

class KnowledgeBaseRetrieveUpdateDestroyView(generics.RetrieveUpdateDestroyAPIView):
    queryset = KnowledgeBase.objects.all()
    serializer_class = KnowledgeBaseSerializer
    permission_classes = [IsAuthenticated]
    lookup_field = 'id'

class SpaceListCreateView(generics.ListCreateAPIView):
    queryset = Space.objects.all().order_by('-created')
    serializer_class = SpaceSerializer
    permission_classes = [IsAuthenticated]
    filter_backends = [DjangoFilterBackend, filters.SearchFilter, filters.OrderingFilter]
    search_fields = ['name']
    ordering_fields = ['created']

class SpaceRetrieveUpdateDestroyView(generics.RetrieveUpdateDestroyAPIView):
    queryset = Space.objects.all()
    serializer_class = SpaceSerializer
    permission_classes = [IsAuthenticated]
    lookup_field = 'id'

class SpaceMemberListCreateView(generics.ListCreateAPIView):
    serializer_class = SpaceMemberSerializer
    permission_classes = [IsAuthenticated]
    def get_queryset(self):
        space_id = self.kwargs['space_id']
        return SpaceMember.objects.filter(space_id=space_id)

class SpaceMemberRetrieveUpdateDestroyView(generics.RetrieveUpdateDestroyAPIView):
    serializer_class = SpaceMemberSerializer
    permission_classes = [IsAuthenticated]
    lookup_field = 'id'
    def get_queryset(self):
        space_id = self.kwargs['space_id']
        return SpaceMember.objects.filter(space_id=space_id)

class SpaceDocumentListCreateView(generics.ListCreateAPIView):
    serializer_class = SpaceDocumentSerializer
    permission_classes = [IsAuthenticated]
    parser_classes = [MultiPartParser, FormParser]
    def get_queryset(self):
        space_id = self.kwargs['space_id']
        return SpaceDocument.objects.filter(space_id=space_id)
    def perform_create(self, serializer):
        file = self.request.FILES.get('file')
        size = f'{round(file.size/1024/1024, 2)}MB' if file else ''
        serializer.save(space_id=self.kwargs['space_id'], size=size)

class SpaceDocumentRetrieveUpdateDestroyView(generics.RetrieveUpdateDestroyAPIView):
    serializer_class = SpaceDocumentSerializer
    permission_classes = [IsAuthenticated]
    lookup_field = 'id'
    def get_queryset(self):
        space_id = self.kwargs['space_id']

def extract_text_from_file(file_path, ext, encoding='utf-8'):
    ext = ext.lower()
    if ext in ['.txt', '.vtt', '.properties']:
        with open(file_path, 'r', encoding=encoding) as f:
            return f.read()
    elif ext in ['.md', '.markdown', '.mdx']:
        with open(file_path, 'r', encoding=encoding) as f:
            return markdown.markdown(f.read())
    elif ext in ['.docx']:
        doc = Document(file_path)
        return '\n'.join([p.text for p in doc.paragraphs])
    elif ext in ['.pdf']:
        with pdfplumber.open(file_path) as pdf:
            return '\n'.join(page.extract_text() or "" for page in pdf.pages)
    elif ext in ['.csv']:
        df = pd.read_csv(file_path)
        return df.to_string(index=False)
    elif ext in ['.xlsx', '.xls']:
        df = pd.read_excel(file_path)
        return df.to_string(index=False)
    elif ext in ['.pptx', '.ppt']:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    elif ext in ['.html', '.htm', '.xml']:
        with open(file_path, 'r', encoding=encoding) as f:
            soup = BeautifulSoup(f, 'html.parser')
            return soup.get_text()
    elif ext in ['.epub']:
        book = epub.read_epub(file_path)
        text = []
        for item in book.get_items():
            if item.get_type() == epub.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                text.append(soup.get_text())
        return '\n'.join(text)
    else:
        return "暂不支持该文件类型"

class KnowledgeFileProcessView(APIView):
    def post(self, request, file_id):
        try:
            params = request.data
            try:
                kf = KnowledgeFile.objects.get(id=file_id)
            except KnowledgeFile.DoesNotExist:
                return Response({"error": "知识文件不存在"}, status=status.HTTP_404_NOT_FOUND)

            # 保存本次分段/embedding参数到知识库
            kb = kf.kb
            kb.embedding_config = params
            kb.save(update_fields=['embedding_config'])

            file_path = kf.file.path
            ext = os.path.splitext(file_path)[-1].lower()
            encoding = params.get("loader_config", {}).get("encoding", "utf-8")
            try:
                file_content = extract_text_from_file(file_path, ext, encoding)
            except Exception as e:
                logger.error(f'[KnowledgeFileProcess] 读取文件内容失败: {str(e)}')
                return Response({"error": f"读取文件内容失败: {str(e)}"}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
            if file_content == "暂不支持该文件类型":
                return Response({"error": file_content}, status=status.HTTP_400_BAD_REQUEST)

            # 2. 文本清洗
            clean_config = params.get('clean_config', {})
            import re

            # 保存原始字符数
            kf.char_count = len(file_content)
            kf.save(update_fields=['char_count'])

            try:
                if clean_config.get('clean_text'):
                    file_content = re.sub(r'[ \t\r\f\v]+', ' ', file_content)
                if clean_config.get('remove_urls'):
                    file_content = re.sub(r'https?://\S+', '', file_content)
                if clean_config.get('remove_emails'):
                    file_content = re.sub(r'\S+@\S+', '', file_content)
                if clean_config.get('remove_extra_whitespace'):
                    file_content = ' '.join(file_content.split())
                if clean_config.get('remove_special_chars'):
                    file_content = re.sub(r'[^\w\s]', '', file_content)
            except Exception as e:
                logger.error(f'[KnowledgeFileProcess] 文本清洗失败: {str(e)}')
                return Response({"error": f"文本清洗失败: {str(e)}"}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

            # 3. 文本分段
            splitter_config = params.get('splitter_config', {})
            split_method = splitter_config.get('text_splitter', 'recursive')
            chunk_size = splitter_config.get('chunk_size', 1000)
            chunk_overlap = splitter_config.get('chunk_overlap', 200)
            separators = splitter_config.get('separators', ['\n\n'])

            # 4. 嵌入生成
            embedding_config = params.get('embedding_config', {})
            embedding_model = embedding_config.get('model', 'sentence-transformers/all-MiniLM-L6-v2')
            embedding_type = embedding_config.get('type', 'local')  # 默认使用本地模型

            user = request.user
            logger.info(f'[KnowledgeFileProcess] 当前用户: {user.username} 开始处理知识库文件')

            # 根据配置选择向量化模型
            try:
                if embedding_type == 'local':
                    # 使用本地模型
                    model_name = embedding_model if embedding_model else "sentence-transformers/all-MiniLM-L6-v2"
                    embedder = HuggingFaceEmbeddings(
                        model_name=model_name,
                        model_kwargs={'device': 'cpu'},
                        encode_kwargs={'normalize_embeddings': True}
                    )
                    logger.info(f'[KnowledgeFileProcess] 使用本地模型 {model_name} 进行向量化')
                elif embedding_type == 'openai':
                    # 使用 OpenAI API
                    api_key = user.openai_api_key
                    from .models import ModelApi
                    model_api = ModelApi.objects.filter(model='openai').order_by('-time').first()
                    api_key = model_api.apikey if model_api else api_key
                    logger.info(f'[KnowledgeFileProcess] 当前用户: {user.username}, OpenAI API密钥: {api_key}')
                    if not api_key:
                        logger.error(f'[KnowledgeFileProcess] 用户 {user.username} 和全局都未设置OpenAI API密钥')
                        return Response({"error": "请先在个人设置或API管理中设置OpenAI API密钥"}, status=status.HTTP_400_BAD_REQUEST)
                    embedder = OpenAIEmbeddings(
                        model=embedding_model,
                        openai_api_key=api_key
                    )
                else:
                    logger.error(f'[KnowledgeFileProcess] 不支持的嵌入模型类型: {embedding_type}')
                    return Response({"error": "不支持的嵌入模型类型"}, status=status.HTTP_400_BAD_REQUEST)
            except Exception as e:
                logger.error(f'[KnowledgeFileProcess] 初始化向量化模型失败: {str(e)}')
                return Response({"error": f"初始化向量化模型失败: {str(e)}"}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

            # 5. 向量存储
            try:
                # 生成嵌入向量
                embeddings = embedder.embed_documents(file_content)
            except Exception as e:
                if 'RateLimitError' in str(e):
                    logger.error(f'[KnowledgeFileProcess] 用户 {user.username} 调用 API 被限流，请稍后重试或更换 API Key')
                    return Response({"error": "API 调用频率超限，请稍后重试或更换 API Key"}, status=status.HTTP_429_TOO_MANY_REQUESTS)
                else:
                    logger.error(f'[KnowledgeFileProcess] 生成嵌入向量失败: {str(e)}')
                    return Response({"error": f"生成嵌入向量失败: {str(e)}"}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

            # 6. 检索索引（Chroma自动支持）
            persist_directory = f"./chroma_db/{file_id}"
            os.makedirs(persist_directory, exist_ok=True)
            vectordb = Chroma.from_texts(
                texts=file_content,
                embedding=embedder,
                persist_directory=persist_directory,
                metadatas=[{"file_id": file_id, "chunk_index": i} for i in range(len(file_content))]
            )
            vectordb.persist()

            # 7. 返回结果
            result = {
                'file_id': file_id,
                'chunk_count': len(file_content),
                'embedding_model': model_name,
                'embedding_type': embedding_type,
                'vector_store_path': persist_directory
            }
            return Response(result)
        except Exception as e:
            logger.error(f'[KnowledgeFileProcess] 未知错误: {str(e)}', exc_info=True)
            return Response({"error": f"未知错误: {str(e)}"}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

class KnowledgeFileUploadView(APIView):
    parser_classes = (MultiPartParser, FormParser)
    permission_classes = [IsAuthenticated]

    def post(self, request, *args, **kwargs):
        logger.info(f'[KnowledgeFileUpload] 开始上传文件，请求数据: {request.data}')
        serializer = KnowledgeFileSerializer(data=request.data)
        if serializer.is_valid():
            try:
                instance = serializer.save()
                logger.info(f'[KnowledgeFileUpload] 文件上传成功: {instance.filename}')
                return Response(serializer.data, status=status.HTTP_201_CREATED)
            except Exception as e:
                logger.error(f'[KnowledgeFileUpload] 保存文件失败: {str(e)}')
                return Response({"error": f"保存文件失败: {str(e)}"}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
        logger.error(f'[KnowledgeFileUpload] 数据验证失败: {serializer.errors}')
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

class KnowledgeFileListView(generics.ListAPIView):
    serializer_class = KnowledgeFileSerializer
    permission_classes = [IsAuthenticated]
    filter_backends = [DjangoFilterBackend, filters.SearchFilter]
    filterset_fields = ['kb']
    search_fields = ['filename']

    def get_queryset(self):
        return KnowledgeFile.objects.all().order_by('-created')